from indic_transliteration import sanscript
from datetime import datetime, timezone
from celery import shared_task
import os, tempfile, platform, zipfile, tarfile, json, yaml, gzip,py7zr, rarfile,filetype
import pandas as pd
import PyPDF2
from docx import Document
import xml.etree.ElementTree as ET
import whisper, ffmpeg
from pptx import Presentation
from pydub import AudioSegment
from PIL import Image
import logging
from .models import (
    AudioFile, VideoFile, ImageFile, DocumentFile, HtmlFile,
    JsonFile, XmlFile, LogFile, PPTFile, SpreadsheetFile, ArchiveFile, YamlFile
)
from .minio_client import get_minio_client, list_objects 
from .utils import detect_file_type, SPREADSHEET_EXTENSIONS, normalize_filename
from .minio_client import move_object
from pathlib import Path
from pymongo import MongoClient
from faster_whisper import WhisperModel
import torch
from bs4 import BeautifulSoup
from celery.signals import worker_process_init
from mongoengine import connect
from xtr.utils import extract_ppt_text

from xtr import minio_client

from .minio_client import get_minio_client

# ✅ Define minio_client for backward compatibility
minio_client = get_minio_client()


###########

@shared_task(bind=True)
def process_minio_file(self, bucket_name, object_name):

    # 🚫 HARD STOP for non-processing buckets (no log)
    if bucket_name != "processing":
        return  # Only process the 'processing' bucket, no log

    logger.warning(
        "[TASK] 🚀 New file event: %s in bucket: %s",
        object_name,
        bucket_name,
    )

    # ✅ Only processing bucket reaches here
    auto_discover_and_process.delay("processing")



@shared_task
def fetch_all_buckets_and_objects():
    print("[TASK] 🚀 Fetching all buckets and their objects from MinIO...")
    skip_bucket = "archive"  # bucket you want to skip
    try:
        for bucket in get_minio_client().list_buckets():
            bucket_name = bucket.name
            if bucket_name == skip_bucket:
                continue  # Skip this bucket!
            print(f"[TASK] 📂 Bucket: {bucket_name}")
            for obj in list_objects(bucket_name):
                filename = obj.object_name
                auto_discover_and_process.delay(bucket_name, filename)
    except Exception as e:
        print(f"[TASK] ❌ Error fetching buckets/objects: {e}")


@shared_task
def auto_discover_and_process(bucket_name=None, filename=None):
    if not bucket_name:
        print("[TASK] ⚠️ Missing bucket_name")
        return
    if bucket_name != "processing":
        return  # Only process the 'processing' bucket

    files = [type("obj", (object,), {"object_name": filename})] if filename else list_objects(bucket_name)

    for obj in files:
        fname = obj.object_name.strip()
        fname = normalize_filename(fname)
        ftype = detect_file_type(fname)
        print(f"[TASK] ➡️ Found: {fname} (type: {ftype})")

        processed = False

        # Dispatch by file type, update processed = True if actually processed
        if ftype == "audio" and not AudioFile.objects(filename=fname).first():
            process_audio.delay(bucket_name, fname)
            processed = True
        elif ftype == "video" and not VideoFile.objects(filename=fname).first():
            process_video.delay(bucket_name, fname)
            processed = True
        elif ftype == "image" and not ImageFile.objects(filename=fname).first():
            process_image.delay(bucket_name, fname)
            processed = True
        elif ftype == "document" and not DocumentFile.objects(filename=fname).first():
            process_doc.delay(bucket_name, fname)
            processed = True
        elif ftype == "presentation" and not PPTFile.objects(filename=fname).first():
            process_ppt.delay(bucket_name, fname)
            processed = True
        elif ftype == "spreadsheet" and not SpreadsheetFile.objects(filename=fname).first():
            process_spreadsheet.delay(bucket_name, fname)
            processed = True
        elif ftype == "html" and not HtmlFile.objects(filename=fname).first():
            process_html.delay(bucket_name, fname)
            processed = True
        elif ftype == "json" and not JsonFile.objects(filename=fname).first():
            process_json.delay(bucket_name, fname)
            processed = True
        elif ftype == "xml" and not XmlFile.objects(filename=fname).first():
            process_xml.delay(bucket_name, fname)
            processed = True
        elif ftype == "log" and not LogFile.objects(filename=fname).first():
            process_log.delay(bucket_name, fname)
            processed = True
        elif ftype == "archive" and not ArchiveFile.objects(filename=fname).first():
            process_archive.delay(bucket_name, fname)
            processed = True
        elif ftype == "yaml" and not YamlFile.objects(filename=fname).first():
            process_yaml.delay(bucket_name, fname)
            processed = True
        else:
            print(f"[TASK] ⏭️ Skipped or unknown: {fname}")



# ----------------------------
# FFMPEG / pydub setup
# ----------------------------
FFMPEG_EXE = os.getenv("FFMPEG_PATH") or "/usr/bin/ffmpeg"
FFPROBE_EXE = os.getenv("FFPROBE_PATH") or "/usr/bin/ffprobe"

AudioSegment.converter = FFMPEG_EXE
AudioSegment.ffmpeg = FFMPEG_EXE
AudioSegment.ffprobe = FFPROBE_EXE

# ----------------------------
# Load Faster-Whisper
# ----------------------------


logger = logging.getLogger(__name__)

_MODEL = None

def get_whisper_model():
    global _MODEL

    if _MODEL is None:
        # Read configuration from environment variables
        model_size = os.getenv("WHISPER_MODEL", "tiny")  # default to 'tiny'
        device_env = os.getenv("WHISPER_DEVICE", None)   # optional override
        compute_type_env = os.getenv("WHISPER_COMPUTE_TYPE", None)  # optional override

        # Determine device: use environment override if provided
        if device_env:
            device = device_env
        else:
            # Auto-detect GPU availability
            device = "cuda" if torch.cuda.is_available() else "cpu"

        # Determine compute_type: use environment override if provided
        if compute_type_env:
            compute_type = compute_type_env
        else:
            compute_type = "float16" if device == "cuda" else "int8"

        try:
            logger.info(f"🚀 Loading Faster-Whisper model '{model_size}' on {device} ({compute_type})...")
            _MODEL = WhisperModel(model_size, device=device, compute_type=compute_type)
            logger.info("✅ Faster-Whisper model loaded successfully.")
        except Exception as e:
            logger.error(f"❌ Failed to load Faster-Whisper model: {e}")
            _MODEL = None

    return _MODEL

# ----------------------------
# Helper: Transcribe file with per-segment logging
# ----------------------------

def transcribe_file(file_path: str, language: str = None):
    if get_whisper_model() is None:
        raise RuntimeError("Faster-Whisper model not loaded")
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    logger.info(f"🔊 Starting transcription for {file_path}")
    segments, info = get_whisper_model().transcribe(file_path, language=language)
    segments = list(segments)
    # Log each segment
    for i, seg in enumerate(segments, start=1):
        logger.info(f"Segment {i}/{len(segments)}: {seg.text.strip()}")

    text = " ".join([seg.text for seg in segments]).strip()
    logger.info(f"✅ Transcription completed: {file_path} | Duration: {info.duration:.2f}s | Language: {info.language}")
    return text, info.language, info.duration

# ----------------------------
# Celery Task: Audio
# ----------------------------
@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def process_audio(self, bucket_name, filename):
    """
    Process audio files from MinIO:
    1. Download from MinIO
    2. Transcribe using Faster-Whisper
    3. Save metadata to MongoDB
    4. Move to archive if successful
    """
    path = None
    status = "failed"
    
    try:
        filename = normalize_filename(filename)
        logger.info(f"[TASK] 🎧 Processing audio: {filename}")

        # Download file from MinIO
        fd, path = tempfile.mkstemp(suffix=os.path.splitext(filename)[-1])

        os.close(fd)

        minio_client.fget_object(bucket_name, filename, path)

        if not os.path.exists(path):
            raise FileNotFoundError(f"Downloaded audio file not found at {path}")

        # Transcribe
        text, detected_lang, duration = transcribe_file(path)

        # Save to MongoDB
        doc = AudioFile(
            filename=filename,
            content=text,
            status="completed",
            meta_data={"detected_language": detected_lang, "duration_sec": duration},
            created_at=datetime.now(timezone.utc)
        )
        doc.save()
        logger.info(f"[TASK] ✅ AudioFile saved with ID: {doc.id}")
        status = "completed"

    except Exception as exc:
        logger.error(f"[TASK] ❌ Failed audio {filename}: {exc}")
        try:
            doc = AudioFile(
                filename=filename,
                content="",
                status="failed",
                meta_data={"error": str(exc)},
                created_at=datetime.now(timezone.utc)
            )
            doc.save()
        except Exception as e:
            logger.error(f"[TASK] ❌ Failed saving failed audio record: {e}")
        try:
            self.retry(exc=exc)
        except self.MaxRetriesExceededError:
            logger.error(f"[TASK] Max retries reached for {filename}")

    finally:
        # Clean up temp file
        if path and os.path.exists(path):
            os.remove(path)

        # Move to archive ONLY if completed
        if bucket_name == "processing" and status == "completed":
            try:
                archive_bucket = "archive"
                success = move_object(bucket_name, filename, archive_bucket)
                if success:
                    logger.info(f"[TASK] 📦 Moved '{filename}' to archive")
                else:
                    logger.warning(f"[TASK] ⏭️ Failed to move '{filename}' to archive")
            except Exception as e:
                logger.error(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")

# ----------------------------
# Celery Task: Video
# ----------------------------
@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def process_video(self, bucket_name, filename):
    vpath, apath = None, None
    try:
        filename = normalize_filename(filename)
        logger.info(f"[TASK] 🎬 Processing video: {filename}")

        # Download video
        fd, vpath = tempfile.mkstemp(suffix=os.path.splitext(filename)[-1])
        os.close(fd)
        minio_client.fget_object(bucket_name, filename, vpath)

        if not os.path.exists(vpath):
            raise FileNotFoundError(f"Downloaded video file not found at {vpath}")

        # Extract audio
        fd, apath = tempfile.mkstemp(suffix=".wav")
        os.close(fd)
        logger.info(f"🎵 Extracting audio from {vpath} to {apath}")

        (
            ffmpeg
            .input(vpath)
            .output(apath, format="wav", acodec="pcm_s16le", ac=1, ar="16000")
            .run(capture_stdout=True, capture_stderr=True, overwrite_output=True)
        )
        logger.info(f"✅ Audio extraction completed: {apath}")

        if not os.path.exists(apath):
            raise FileNotFoundError(f"Extracted audio file not found at {apath}")

        # Transcribe audio
        text, detected_lang, duration = transcribe_file(apath)

        doc = VideoFile(
            filename=filename,
            content=text,
            status="completed",
            meta_data={"detected_language": detected_lang, "duration_sec": duration},
            created_at=datetime.now(timezone.utc)
        )
        doc.save()
        logger.info(f"[TASK] ✅ VideoFile saved with ID: {doc.id}")

    except Exception as exc:
        logger.error(f"[TASK] ❌ Failed video {filename}: {exc}")
        try:
            doc = VideoFile(
                filename=filename,
                content="",
                status="failed",
                meta_data={"error": str(exc)},
                created_at=datetime.now(timezone.utc)
            )
            doc.save()
        except Exception as e:
            logger.error(f"[TASK] ❌ Failed saving failed video record: {e}")
        try:
            self.retry(exc=exc)
        except self.MaxRetriesExceededError:
            logger.error(f"[TASK] Max retries reached for {filename}")

    finally:
        for p in [vpath, apath]:
            if p and os.path.exists(p):
                os.remove(p)
        if bucket_name == "processing":
            try:
                archive_bucket = "archive"
                doc = VideoFile.objects(filename=filename).first()
                status = doc.status if doc else "completed"
                success = move_object(bucket_name, filename, archive_bucket)
                if success:
                    logger.info(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' to '{archive_bucket}'")
            except Exception as e:
                logger.error(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")

# ------------------------------------
# IMAGE TASK
# ------------------------------------


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def process_image(self, bucket_name, object_name):
    """
    Task: Process image files from MinIO 'processing' bucket.
    Extracts metadata and moves completed files to 'archive' bucket.
    """
    tmp_path = None
    status = "failed"  # default status, will change to completed later

    try:
        object_name = normalize_filename(object_name)
        logger.info(f"[TASK] 📷 Processing image: {object_name}")

        # 1️⃣ Download file from MinIO
        ext = os.path.splitext(object_name)[-1].lower()
        fd, tmp_path = tempfile.mkstemp(suffix=ext)
        os.close(fd)
        minio_client.fget_object(bucket_name, object_name, tmp_path)

        if not os.path.exists(tmp_path):
            raise FileNotFoundError(f"Downloaded image file not found at {tmp_path}")

        # 2️⃣ Handle extension typos (common mistranscriptions from uploads/events)
        # Added .jif -> .jfif and a few common misspellings to match IMAGE_EXTENSIONS
        typo_map = {
            ".ppng": ".png",
            ".jiif": ".jfif",
            ".jif": ".jfif",
            ".jgp": ".jpg",
            ".jpe": ".jpeg",
        }
        if ext in typo_map:
            new_path = tmp_path.replace(ext, typo_map[ext])
            os.rename(tmp_path, new_path)
            tmp_path = new_path
            ext = typo_map[ext]

        # 3️⃣ Enable HEIC/HEIF support
        if ext in [".heic", ".heif"]:
            try:
                from pillow_heif import register_heif_opener
                register_heif_opener()
            except Exception as e:
                logger.warning(f"[TASK] ⚠️ Could not load HEIC support: {e}")

        # 4️⃣ Handle SVG (convert to PNG first)
        if ext == ".svg":
            try:
                import cairosvg
                png_path = tmp_path + ".png"
                cairosvg.svg2png(url=tmp_path, write_to=png_path)
                os.remove(tmp_path)
                tmp_path = png_path
                ext = ".png"
            except Exception as e:
                logger.error(f"[TASK] ❌ Failed to convert SVG: {e}")
                raise

        # 5️⃣ Open with Pillow and extract metadata
        with Image.open(tmp_path) as img:
            # Ensure standard mode
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")

            # Clean metadata
            clean_info = {}
            for k, v in img.info.items():
                try:
                    clean_info[k] = float(v) if hasattr(v, "__float__") else str(v)
                except Exception:
                    clean_info[k] = str(v)

            # --- Upsert instead of insert ---
            ImageFile.objects(filename=object_name).update_one(
                set__file_size=os.path.getsize(tmp_path),
                set__width=img.size[0],
                set__height=img.size[1],
                set__format=img.format or ext,
                set__status="completed",
                set__meta_data={"mode": img.mode, "info": clean_info},
                set__created_at=datetime.now(timezone.utc),
                upsert=True
            )
            logger.info(f"[TASK] ✅ ImageFile saved: {object_name}")

        # ✅ Mark as completed
        status = "completed"

    except Exception as exc:
        logger.error(f"[TASK] ❌ Failed image {object_name}: {exc}")
        try:
            ImageFile.objects(filename=object_name).update_one(
                set__status="failed",
                set__meta_data={"error": str(exc)},
                set__created_at=datetime.now(timezone.utc),
                upsert=True
            )
        except Exception as e:
            logger.error(f"[TASK] ❌ Failed saving failed image record: {e}")
        try:
            self.retry(exc=exc)
        except self.MaxRetriesExceededError:
            logger.error(f"[TASK] Max retries reached for {object_name}")

    finally:
        # ✅ Clean up temp file
        if tmp_path and os.path.exists(tmp_path):
            os.remove(tmp_path)
        # ✅ Move image to archive similar to other tasks when processed from 'processing'
        if bucket_name == "processing":
            try:
                # determine status from DB record if present
                img_doc = ImageFile.objects(filename=object_name).first()
                status = getattr(img_doc, 'status', 'completed') if img_doc else 'completed'
                archive_bucket = os.environ.get('MINIO_ARCHIVE_BUCKET', 'archive')
                moved = move_object(bucket_name, object_name, archive_bucket)
                if moved:
                    logger.info("[TASK] 📦 Moved image '%s' to '%s' (status=%s)", object_name, archive_bucket, status)
                else:
                    logger.info("[TASK] ⏭️ Image '%s' not moved (status=%s)", object_name, status)
            except Exception as e:
                logger.error("[TASK] ⚠️ Could not move image '%s' to archive: %s", object_name, e)



# ------------------------------------
# DOCUMENT TASK
# ------------------------------------


@shared_task
def process_doc(bucket_name, object_name):
    print(f"[TASK] 📄 Document: {object_name}")
    ext = os.path.splitext(object_name)[-1].lower()
    tmp = None
    try:
        fd, tmp = tempfile.mkstemp(suffix=ext)
        os.close(fd)
        object_name=normalize_filename(object_name)
        minio_client.fget_object(bucket_name, object_name, tmp)

        text = ""

        if ext == ".pdf":
            with open(tmp, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = "\n".join(p.extract_text() or "" for p in reader.pages)

        elif ext == ".docx":
            doc = Document(tmp)
            text = "\n".join(p.text for p in doc.paragraphs)

        elif ext == ".odt":
            from odf.opendocument import load
            from odf import text as odf_text
            odt_doc = load(tmp)
            parts = []
            for elem in odt_doc.getElementsByType(odf_text.P):
                parts.append(str(elem))
            text = "\n".join(parts)

        elif ext == ".epub":
            from ebooklib import epub
            from bs4 import BeautifulSoup
            book = epub.read_epub(tmp)
            parts = []
            for item in book.get_items_of_type(9):  # DOCUMENT
                soup = BeautifulSoup(item.get_content(), "html.parser")
                parts.append(soup.get_text(" ", strip=True))
            text = "\n".join(parts)

        else:
            # fallback for .txt or unknown
            with open(tmp, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        DocumentFile.objects.create(
            filename=object_name,
            content=text,
            status="completed",
            meta_data={"ext": ext, "length": len(text)},
        )
        print(f"[TASK] ✅ Document processed: {object_name}")

    except Exception as e:
        DocumentFile.objects.create(
            filename=object_name,
            content="",
            status="failed",
            meta_data={"error": str(e)},
        )
        print(f"[TASK] ❌ Failed processing {object_name}: {e}")

    finally:
        if tmp and os.path.exists(tmp):
            os.remove(tmp)

        if bucket_name == "processing":
            try:
                archive_bucket = "archive"
                doc = DocumentFile.objects(filename=object_name).first()
                status = doc.status if doc else "completed"
                success = move_object(bucket_name, object_name, archive_bucket)
                if success:
                    print(f"[TASK] 📦 Moved '{object_name}' from '{bucket_name}' to '{archive_bucket}'")
            except Exception as e:
                print(f"[TASK] ⚠️ Could not move '{object_name}' to archive: {e}")

# ------------------------------------
# HTML TASK
# ------------------------------------



@shared_task
def process_html(bucket_name, filename):
    try:
        # Fetch HTML from MinIO
        data = minio_client.get_object(bucket_name, filename).read().decode()

        # Parse HTML bookmarks
        soup = BeautifulSoup(data, "html.parser")
        bookmarks = []

        for a_tag in soup.find_all('a'):
            bookmarks.append({
                "name": a_tag.get_text(),
                "url": a_tag.get('href')
            })

        # Convert to human-readable string (or JSON if you want)
        human_readable = "\n".join([f"{b['name']}: {b['url']}" for b in bookmarks])

        # Save in database
        HtmlFile.objects.create(
            filename=filename,
            content=human_readable,
            status="completed"
        )
        print(f"[TASK] ✅ HTML processed and converted: {filename}")

    except Exception as e:
        HtmlFile.objects.create(
            filename=filename,
            content="",
            status="failed",
            meta_data={"error": str(e)}
        )
    if bucket_name == "processing":
        try:
            archive_bucket = "archive"
            doc = HtmlFile.objects(filename=filename).first()
            status = doc.status if doc else "completed"
            success = move_object(bucket_name, filename, archive_bucket)
            if success:
                print(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' to '{archive_bucket}'")
        except Exception as e:
            print(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")

# ------------------------------------
# JSON TASK
# ------------------------------------


@shared_task
def process_json(bucket_name, filename):
    try:
        raw = minio_client.get_object(bucket_name, filename).read().decode()
        JsonFile.objects.create(filename=filename, content=raw, status="completed")
        print(f"[TASK] ✅ JSON processed: {filename}")
    except Exception as e:
        JsonFile.objects.create(filename=filename, content="", status="failed", meta_data={"error": str(e)})
    if bucket_name == "processing":
        try:
            archive_bucket = "archive"
            doc = JsonFile.objects(filename=filename).first()
            status = doc.status if doc else "completed"
            success = move_object(bucket_name, filename, archive_bucket)
            if success:
                print(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' to '{archive_bucket}'")
        except Exception as e:
            print(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")

# ------------------------------------
# XML TASK
# ------------------------------------


@shared_task
def process_xml(bucket_name, filename):
    try:
        raw = minio_client.get_object(bucket_name, filename).read().decode()
        root = ET.fromstring(raw)
        XmlFile.objects.create(filename=filename, content=raw, status="completed", meta_data={"root_tag": root.tag})
        print(f"[TASK] ✅ XML processed: {filename}")
    except Exception as e:
        XmlFile.objects.create(filename=filename, content="", status="failed", meta_data={"error": str(e)})
    if bucket_name == "processing":
        archive_bucket = "archive"
        status = "completed"  # Set status based on processing outcome
        success = move_object(bucket_name, filename, archive_bucket, status)
        if success:
            print(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' to '{archive_bucket}'")
        else:
            print(f"[TASK] ⚠️ Could not move '{filename}' to archive.")

# ------------------------------------
# LOG TASK
# ------------------------------------

@shared_task
def process_log(bucket_name, filename):
    try:
        raw = minio_client.get_object(bucket_name, filename).read().decode()
        LogFile.objects.create(filename=filename, content=raw, status="completed", meta_data={"length": len(raw)})
        print(f"[TASK] ✅ Log processed: {filename}")
    except Exception as e:
        LogFile.objects.create(filename=filename, content="", status="failed", meta_data={"error": str(e)})
    if bucket_name == "processing":
        try:
            archive_bucket = "archive"
            doc = LogFile.objects(filename=filename).first()
            status = doc.status if doc else "completed"
            success = move_object(bucket_name, filename, archive_bucket)
            if success:
                print(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' to '{archive_bucket}'")
        except Exception as e:
            print(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")

# ------------------------------------
#  YAML TASK
# ------------------------------------

@shared_task(bind=True)
def process_yaml(self, bucket_name, filename):
    tmp_file = None
    try:
        print(f"[TASK] 📄 YAML: {filename}")

        # Create temp file safely
        fd, tmp_file = tempfile.mkstemp(suffix=".yaml")
        os.close(fd)

        # Download from MinIO
        minio_client.fget_object(bucket_name, filename, tmp_file)

        # Read YAML safely
        with open(tmp_file, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)

        # Save to MongoDB
        YamlFile.objects.create(
            filename=filename,
            content=yaml.dump(data),
            status="completed",
            meta_data={"keys": list(data.keys()) if isinstance(data, dict) else None},
        )

        print(f"[TASK] ✅ YAML processed: {filename}")

    except Exception as e:
        # Save failure info
        YamlFile.objects.create(
            filename=filename,
            content="",
            status="failed",
            meta_data={"error": str(e)},
        )
        print(f"[TASK] ❌ Failed processing {filename}: {e}")

    finally:
        # Clean up temp file
        if tmp_file and os.path.exists(tmp_file):
            try:
                os.remove(tmp_file)
            except Exception as cleanup_error:
                print(f"[WARNING] Could not remove temp file {tmp_file}: {cleanup_error}")

        if bucket_name == "processing":
            try:
                archive_bucket = "archive"
                doc = YamlFile.objects(filename=filename).first()
                status = doc.status if doc else "completed"
                success = move_object(bucket_name, filename, archive_bucket)
                if success:
                    print(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' to '{archive_bucket}'")
            except Exception as e:
                print(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")
                
# ------------------------------------
# PPT TASK
# ------------------------------------

logger = logging.getLogger(__name__)


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def process_ppt(self, bucket_name, filename):
    """
    Task: Process PPT/PPTX files from MinIO 'processing' bucket.
    Extracts text content and moves completed files to 'archive' bucket.
    """
    path = None
    status = "failed"  # default status, will change to completed later

    try:
        filename = normalize_filename(filename)
        logger.info(f"[TASK] 📊 Processing PPT file: {filename}")

        # Temporary local file
        fd, path = tempfile.mkstemp(suffix=os.path.splitext(filename)[-1])
        os.close(fd)

        # Download file from MinIO
        minio_client.fget_object(bucket_name, filename, path)

        if not os.path.exists(path):
            raise FileNotFoundError(f"Downloaded file not found: {path}")

        # Extract PPT text and slide count
        extracted_text, slide_count = extract_ppt_text(path)

        # Save document info to MongoDB
        doc = PPTFile(
            filename=filename,
            content=extracted_text,
            status="completed",
            meta_data={"slides": slide_count},
            created_at=datetime.now(timezone.utc)
        )
        doc.save()
        status = "completed"
        logger.info(f"[TASK] ✅ PPTFile saved successfully: {doc.id}")

    except Exception as exc:
        logger.error(f"[TASK] ❌ Failed to process {filename}: {exc}")
        try:
            doc = PPTFile(
                filename=filename,
                content="",
                status="failed",
                meta_data={"error": str(exc)},
                created_at=datetime.now(timezone.utc)
            )
            doc.save()
        except Exception as e:
            logger.error(f"[TASK] ❌ Failed saving failed PPT record: {e}")

        try:
            self.retry(exc=exc)
        except self.MaxRetriesExceededError:
            logger.error(f"[TASK] Max retries reached for {filename}")

    finally:
        # ✅ Move to archive only if status == "completed"
        if bucket_name == "processing" and status == "completed":
            try:
                archive_bucket = "archive"
                success = move_object(bucket_name, filename, archive_bucket)
                if success:
                    logger.info(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' → '{archive_bucket}'")
            except Exception as e:
                logger.error(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")

        # 🧹 Clean up temporary file
        if path and os.path.exists(path):
            os.remove(path)
            logger.info(f"[TASK] 🧹 Cleaned up temp file: {path}")

# @shared_task
# def process_ppt(bucket_name, filename):
#     tmp, converted = None, None
#     try:
#         ext = os.path.splitext(filename)[-1].lower()

#         # Download file from MinIO
#         fd, tmp = tempfile.mkstemp(suffix=ext)
#         os.close(fd)
#         minio_client.fget_object(bucket_name, filename, tmp)

#         # Handle .ppt (convert to .pptx)
#         if ext == ".ppt":
#             outdir = os.path.dirname(tmp)
#             basename = os.path.splitext(os.path.basename(tmp))[0]
#             libreoffice_bin = os.environ.get("LIBREOFFICE_BIN", "libreoffice")
#             try:
#                 subprocess.run(
#                     [libreoffice_bin, "--headless", "--convert-to", "pptx", "--outdir", outdir, tmp],
#                     check=True,
#                     stdout=subprocess.PIPE,
#                     stderr=subprocess.PIPE,
#                 )
#             except subprocess.CalledProcessError as e:
#                 # Try fallback to 'soffice' binary
#                 fallback_bin = os.environ.get("LIBREOFFICE_FALLBACK_BIN", "soffice")
#                 subprocess.run(
#                     [fallback_bin, "--headless", "--convert-to", "pptx", "--outdir", outdir, tmp],
#                     check=True,
#                     stdout=subprocess.PIPE,
#                     stderr=subprocess.PIPE,
#                 )
#             converted = os.path.join(outdir, basename + ".pptx")
#             if not os.path.exists(converted):
#                 raise RuntimeError("Failed to convert .ppt to .pptx. Ensure LibreOffice is installed.")
#             tmp = converted  # switch to converted file

#         # Read presentation
#         prs = Presentation(tmp)
#         text = "\n".join(
#             [sh.text for sl in prs.slides for sh in sl.shapes if hasattr(sh, "text")]
#         )

#         # Save result in DB
#         PPTFile.objects.create(
#             filename=filename,
#             content=text,
#             status="completed",
#             meta_data={"slides": len(prs.slides)},
#         )
#         print(f"[TASK] ✅ PPT processed: {filename}")

#     except Exception as e:
#         PPTFile.objects.create(
#             filename=filename,
#             content="",
#             status="failed",
#             meta_data={"error": str(e)},
#         )
#         print(f"[TASK] ❌ Failed PPT {filename}: {e}")

#     finally:
#         for f in [tmp, converted]:
#             if f and os.path.exists(f):
#                 try:
#                     os.remove(f)
#                 except Exception:
#                     pass
#             if bucket_name == "processing":
#                 try:
#                     from xtr.utils import move_file_to_archive  # adjust import if needed
#                     archive_bucket = "archive"
#                     move_file_to_archive(minio_client, bucket_name, filename, archive_bucket, status="completed")
#                     print(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' to '{archive_bucket}'")
#                 except Exception as e:
#                     print(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")

# -------------------------------------
# SPREADSHEET TASK
# -------------------------------------


@shared_task(bind=True, autoretry_for=(OSError,), retry_backoff=True, max_retries=3)
def process_spreadsheet(self, bucket_name, filename):
    tmp = None
    try:
        ext = os.path.splitext(filename)[-1].lower()
        fd, tmp = tempfile.mkstemp(suffix=ext)
        os.close(fd)

        # Download file from MinIO
        minio_client.fget_object(bucket_name, filename, tmp)

        # Read spreadsheet based on extension
        if ext == ".csv":
            df = pd.read_csv(tmp)

        elif ext == ".xlsx":
            try:
                import openpyxl
                df = pd.read_excel(tmp, engine="openpyxl")
            except ImportError as ie:
                raise RuntimeError(
                    "Missing dependency 'openpyxl'. Install with `pip install openpyxl`."
                ) from ie

        elif ext == ".xls":
            try:
                import xlrd
                df = pd.read_excel(tmp, engine="xlrd")
            except ImportError as ie:
                raise RuntimeError(
                    "Missing dependency 'xlrd'. Install with `pip install xlrd`."
                ) from ie

        elif ext == ".ods":
            try:
                df = pd.read_excel(tmp, engine="odf")
            except ImportError as ie:
                raise RuntimeError(
                    "Missing dependency 'odfpy'. Install with `pip install odfpy`."
                ) from ie

        else:
            raise ValueError(f"Unsupported spreadsheet format: {ext}")

        # Save result
        SpreadsheetFile.objects.create(
            filename=filename,
            content=df.to_csv(index=False),
            status="completed",
            meta_data={"columns": list(df.columns), "num_rows": len(df)},
        )
        print(f"[TASK] ✅ Spreadsheet processed: {filename}")

    except Exception as e:
        SpreadsheetFile.objects.create(
            filename=filename,
            content="",
            status="failed",
            meta_data={"error": str(e)},
        )
        print(f"[TASK] ❌ Failed spreadsheet {filename}: {e}")

    finally:
        if tmp and os.path.exists(tmp):
            os.remove(tmp)

        if bucket_name == "processing":
            try:
                archive_bucket = "archive"
                doc = SpreadsheetFile.objects(filename=filename).first()
                status = doc.status if doc else "completed"
                success = move_object(bucket_name, filename, archive_bucket)
                if success:
                    print(f"[TASK] 📦 Moved '{filename}' from '{bucket_name}' to '{archive_bucket}'")
            except Exception as e:
                print(f"[TASK] ⚠️ Could not move '{filename}' to archive: {e}")

# ------------------------------------
# ARCHIVE TASK
# ------------------------------------



@shared_task
def process_archive(bucket_name, object_name):
    print(f"[TASK] ➡️ Processing archive: {object_name}")
    ext = os.path.splitext(object_name)[-1].lower()
    tmp = None

    try:
        # Download object from MinIO
        fd, tmp = tempfile.mkstemp(suffix=ext)
        os.close(fd)
        minio_client.fget_object(bucket_name, object_name, tmp)

        file_list = []

        # -------- Detect real filetype --------
        kind = filetype.guess(tmp)
        if kind:
            print(f"[DEBUG] Detected type: {kind.mime} ({kind.extension})")
            real_ext = f".{kind.extension}"
        else:
            print("[DEBUG] Could not detect file type, falling back to extension")
            real_ext = ext

        # -------- ZIP --------
        if real_ext == ".zip":
            print("[DEBUG] Using zipfile")
            with zipfile.ZipFile(tmp, "r") as zf:
                file_list = zf.namelist()

        # -------- TAR & compressed TAR --------
        elif real_ext in (".tar", ".tgz", ".tar.gz", ".tar.bz2", ".tar.xz"):
            print("[DEBUG] Using tarfile (auto-detect compression)")
            with tarfile.open(tmp, "r:*") as tf:
                file_list = tf.getnames()

        # -------- GZ (single file, not tar.gz) --------
        elif real_ext == ".gz":
            print("[DEBUG] Using gzip (single file)")
            try:
                with gzip.open(tmp, "rb") as gz:
                    inner_name = os.path.basename(object_name).replace(".gz", "")
                    out_tmp = tmp + "_decompressed"
                    with open(out_tmp, "wb") as out_f:
                        out_f.write(gz.read())
                    file_list = [inner_name]
            except OSError:
                raise RuntimeError("File has .gz extension but is not a valid gzip file")

        # -------- 7Z --------
        elif real_ext == ".7z":
            print("[DEBUG] Using py7zr")
            with open(tmp, "rb") as f:
                sig = f.read(6)
            if sig != b"7z\xbc\xaf\x27\x1c":
                raise RuntimeError("Not a valid 7z archive (wrong header)")
            with py7zr.SevenZipFile(tmp, "r") as zf:
                file_list = zf.getnames()

        # -------- RAR --------
        elif real_ext == ".rar":
            print("[DEBUG] Using rarfile")
            try:
                with rarfile.RarFile(tmp, "r") as rf:
                    file_list = rf.namelist()
            except rarfile.Error as e:
                raise RuntimeError(f"Invalid RAR file (unrar not installed?): {e}")

        else:
            raise RuntimeError(f"Unsupported or unrecognized archive type: {real_ext}")

        # Save success
        ArchiveFile.objects.create(
            filename=object_name,
            content="\n".join(file_list),
            status="completed",
            meta_data={"num_files": len(file_list)}
        )
        print(f"[TASK] ✅ Archive processed: {object_name} with {len(file_list)} files")

    except Exception as e:
        # Save failure
        ArchiveFile.objects.create(
            filename=object_name,
            content="",
            status="failed",
            meta_data={"error": str(e)}
        )
        print(f"[TASK] ❌ Failed processing {object_name}: {e}")

    finally:
        if tmp and os.path.exists(tmp):
            try:
                os.remove(tmp)
            except Exception:
                pass

        if bucket_name == "processing":
            try:
                archive_bucket = "archive"
                doc = ArchiveFile.objects(filename=object_name).first()
                status = doc.status if doc else "completed"
                success = move_object(bucket_name, object_name, archive_bucket)
                if success:
                    print(f"[TASK] 📦 Moved '{object_name}' from '{bucket_name}' to '{archive_bucket}'")
            except Exception as e:
                print(f"[TASK] ⚠️ Could not move '{object_name}' to archive: {e}")


