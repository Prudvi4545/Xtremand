import os
import logging
from pptx import Presentation

from minio.error import S3Error
from minio.commonconfig import CopySource
from xtr.minio_client import get_minio_client

logger = logging.getLogger(__name__)

# ✅ AUDIO
AUDIO_EXTENSIONS = frozenset({
    '.mp3', '.wav', '.aac', '.flac', '.m4a', '.ogg',
    '.wma', '.alac', '.aiff', '.amr', '.opus', '.au',
    '.mp2',  '.mka',  '.voc',  '.ra'
})


# ✅ VIDEO
VIDEO_EXTENSIONS = frozenset({
    '.mp4', '.mkv', '.mov', '.avi', '.webm', '.m4v', '.flv', '.wmv',
    '.mpeg', '.mpg',  '.3gp', '.mxf',  '.vob',
})

# ✅ IMAGE
IMAGE_EXTENSIONS = frozenset({
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff',
    '.heic', '.webp', '.raw', '.svg', '.ico',
    '.psd', '.ai', '.indd', '.jfif',
})

# ✅ DOCUMENT
DOCUMENT_EXTENSIONS = frozenset({
    '.txt', '.doc', '.docx', '.rtf', '.odt', '.epub', '.pdf', '.md'
})

# ✅ PRESENTATION
PRESENTATION_EXTENSIONS = frozenset({
    '.ppt', '.pptx'
})

# ✅ SPREADSHEET
SPREADSHEET_EXTENSIONS = frozenset({
    '.csv', '.xlsx', '.xls', '.ods'
})

# ✅ HTML
HTML_EXTENSIONS = frozenset({
    '.html', '.htm'
})

# ✅ SEMI-STRUCTURED
JSON_EXTENSIONS = frozenset({'.json'})
XML_EXTENSIONS = frozenset({'.xml'})
YAML_EXTENSIONS = frozenset({'.yaml', '.yml'})
LOG_EXTENSIONS = frozenset({'.log'})

# ✅ ARCHIVE
ARCHIVE_EXTENSIONS = frozenset({'.zip', '.tar', '.gz', '.rar', '.7z'})


def detect_file_type(filename, debug=False):
    filename = filename.strip()
    ext = os.path.splitext(filename)[1].lower().strip()

    # Fallback: check for double extensions
    if ext in ['.bak', '.tmp', '.old']:
        base = os.path.splitext(filename.strip())[0]
        ext = os.path.splitext(base)[1].lower().strip()

    if debug:
        print(f"🔍 EXT: '{ext}' from filename: '{filename}'")

    if ext in AUDIO_EXTENSIONS:
        return 'audio'
    elif ext in VIDEO_EXTENSIONS:
        return 'video'
    elif ext in IMAGE_EXTENSIONS:
        return 'image'
    elif ext in DOCUMENT_EXTENSIONS:
        return 'document'
    elif ext in PRESENTATION_EXTENSIONS:
        return 'presentation'
    elif ext in SPREADSHEET_EXTENSIONS:
        return 'spreadsheet'
    elif ext in HTML_EXTENSIONS:
        return 'html'
    elif ext in JSON_EXTENSIONS:
        return 'json'
    elif ext in XML_EXTENSIONS:
        return 'xml'
    elif ext in LOG_EXTENSIONS:
        return 'log'
    elif ext in ARCHIVE_EXTENSIONS:
        return 'archive'
    elif ext in YAML_EXTENSIONS:
        return 'yaml'
    else:
        return 'unknown'


def extract_ppt_text(file_path):
    prs = Presentation(file_path)
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides_text.append(shape.text)
    return "\n".join(slides_text), len(prs.slides)


def extract_text_file(file_path):
    with open(file_path, encoding='utf-8', errors='ignore') as f:
        return f.read()


if __name__ == "__main__":
    files = [
        "song.mp3", "video.mp4", "pic.jpg", "doc.pdf", "report.docx",
        "slides.pptx", "data.csv", "page.html", "notes.txt", "archive.zip",
        "weird_file.mp3.bak", "upper.MP3"
    ]
    for f in files:
        print(f"{f}: {detect_file_type(f, debug=True)}")


# utils.py
import urllib.parse

def normalize_filename(filename: str) -> str:
    """
    Fix issues like spaces, %20, + signs from MinIO event filenames.
    Example: "How+to+Copy%20Dashboard.mp4" → "How to Copy Dashboard.mp4"
    """
    return urllib.parse.unquote(filename.replace("+", " "))






def move_file_to_archive(minio_client, source_bucket, object_key, archive_bucket, status):
    """
    Moves a file from processing → archive only when status == 'completed'.
    """
    if status != "completed":
        logger.debug("Skipping move: %s not completed (status=%s)", object_key, status)
        return False

    try:
        client = get_minio_client()
        source = CopySource(source_bucket, object_key)
        client.copy_object(archive_bucket, object_key, source)
        logger.info("Copied '%s' from '%s' to '%s'", object_key, source_bucket, archive_bucket)
        client.remove_object(source_bucket, object_key)
        logger.info("Deleted '%s' from '%s'", object_key, source_bucket)
        return True
    except S3Error as e:
        logger.error("S3Error moving '%s' to '%s': %s", object_key, archive_bucket, e)
        return False
    except Exception as e:
        logger.exception("Unexpected error moving '%s' to '%s': %s", object_key, archive_bucket, e)
        return False